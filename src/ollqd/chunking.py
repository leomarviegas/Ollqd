"""Code-aware and document chunking strategies."""

from pathlib import Path
from typing import Optional

from .models import Chunk, FileInfo


def chunk_pdf(
    file_path: str,
    pdf_bytes: bytes,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    content_hash: str = "",
) -> list[Chunk]:
    """Extract text from PDF and chunk by paragraph boundaries."""
    import fitz  # PyMuPDF

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages_text = []
    for page in doc:
        pages_text.append(page.get_text("text"))
    doc.close()

    full_text = "\n\n".join(pages_text)
    if not full_text.strip():
        return []

    return chunk_document(
        file_path=file_path,
        content=full_text,
        language="text",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        content_hash=content_hash,
    )


def chunk_docx(
    file_path: str,
    docx_bytes: bytes,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    content_hash: str = "",
) -> list[Chunk]:
    """Extract text from .docx (paragraphs + tables) and chunk."""
    from io import BytesIO
    from docx import Document

    doc = Document(BytesIO(docx_bytes))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    full_text = "\n\n".join(parts)
    if not full_text.strip():
        return []

    return chunk_document(
        file_path=file_path,
        content=full_text,
        language="text",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        content_hash=content_hash,
    )


def chunk_xlsx(
    file_path: str,
    xlsx_bytes: bytes,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    content_hash: str = "",
) -> list[Chunk]:
    """Extract text from .xlsx (all sheets, rows as pipe-delimited) and chunk."""
    from io import BytesIO
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(xlsx_bytes), read_only=True, data_only=True)
    sheets_text: list[str] = []
    for ws in wb:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets_text.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
    wb.close()

    full_text = "\n\n".join(sheets_text)
    if not full_text.strip():
        return []

    return chunk_document(
        file_path=file_path,
        content=full_text,
        language="text",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        content_hash=content_hash,
    )


def chunk_pptx(
    file_path: str,
    pptx_bytes: bytes,
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    content_hash: str = "",
) -> list[Chunk]:
    """Extract text from .pptx (slide text frames) and chunk."""
    from io import BytesIO
    from pptx import Presentation

    prs = Presentation(BytesIO(pptx_bytes))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"Slide {i}\n" + "\n".join(texts))

    full_text = "\n\n".join(slides_text)
    if not full_text.strip():
        return []

    return chunk_document(
        file_path=file_path,
        content=full_text,
        language="text",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        content_hash=content_hash,
    )


def _is_boundary_line(line: str, language: str) -> bool:
    """Heuristic: is this line a natural split point?"""
    stripped = line.strip()
    if not stripped:
        return False

    # Markdown headings start with # — check before the comment guard
    if language == "markdown":
        return stripped.startswith(("#", "##", "###"))

    if stripped.startswith("#") or stripped.startswith("//"):
        return False

    if language == "python":
        return stripped.startswith(("def ", "class ", "async def ", "@"))
    if language == "go":
        return stripped.startswith("func ") or stripped.startswith("type ")
    if language in ("javascript", "typescript"):
        return any(stripped.startswith(kw) for kw in (
            "function ", "export ", "class ", "const ", "async function",
            "describe(", "it(", "test(",
        ))
    if language == "rust":
        return any(stripped.startswith(kw) for kw in (
            "fn ", "pub fn ", "impl ", "struct ", "enum ", "mod ", "trait ",
        ))
    if language in ("java", "kotlin", "csharp", "scala"):
        return any(stripped.startswith(kw) for kw in (
            "public ", "private ", "protected ", "class ", "interface ",
            "fun ", "data class ", "object ", "override ",
        ))
    if language in ("c", "cpp"):
        return ("(" in stripped and ")" in stripped and "{" in stripped
                and not stripped.startswith("if")
                and not stripped.startswith("for")
                and not stripped.startswith("while"))
    return False


def chunk_file(file_info: FileInfo, chunk_size: int = 512, chunk_overlap: int = 64) -> list[Chunk]:
    """Split a file into overlapping chunks, preferring natural code boundaries."""
    try:
        content = Path(file_info.abs_path).read_text(errors="replace")
    except (OSError, PermissionError):
        return []

    lines = content.splitlines(keepends=True)
    if not lines:
        return []

    char_budget = chunk_size * 4
    overlap_chars = chunk_overlap * 4

    chunks: list[Chunk] = []
    current_lines: list[str] = []
    current_chars = 0
    chunk_start_line = 1

    def _flush(end_line: int):
        text = "".join(current_lines).strip()
        if text:
            chunks.append(Chunk(
                file_path=file_info.path,
                language=file_info.language,
                chunk_index=len(chunks),
                total_chunks=-1,
                start_line=chunk_start_line,
                end_line=end_line,
                content=text,
                content_hash=file_info.content_hash,
            ))

    for i, line in enumerate(lines, start=1):
        line_len = len(line)

        if (current_chars + line_len > char_budget
                and current_chars > overlap_chars
                and _is_boundary_line(line, file_info.language)):
            _flush(i - 1)
            overlap_text = "".join(current_lines)
            if len(overlap_text) > overlap_chars:
                overlap_text = overlap_text[-overlap_chars:]
            overlap_lines = overlap_text.splitlines(keepends=True)
            current_lines = overlap_lines
            current_chars = sum(len(ln) for ln in current_lines)
            chunk_start_line = max(1, i - len(overlap_lines))

        if current_chars + line_len > char_budget * 1.5 and current_chars > 0:
            _flush(i - 1)
            current_lines = []
            current_chars = 0
            chunk_start_line = i

        current_lines.append(line)
        current_chars += line_len

    _flush(len(lines))

    for c in chunks:
        c.total_chunks = len(chunks)

    return chunks


def chunk_document(
    file_path: str,
    content: str,
    language: str = "text",
    chunk_size: int = 512,
    chunk_overlap: int = 64,
    content_hash: str = "",
) -> list[Chunk]:
    """Chunk a document (markdown, text, etc.) by paragraph boundaries."""
    import re

    if language == "markdown":
        sections = re.split(r"(?=^#{1,3}\s)", content, flags=re.MULTILINE)
    else:
        sections = re.split(r"\n\s*\n", content)

    char_budget = chunk_size
    chunks: list[Chunk] = []
    current_text = ""
    current_start = 1

    for section in sections:
        section = section.strip()
        if not section:
            continue

        if len(current_text) + len(section) > char_budget and current_text:
            chunks.append(Chunk(
                file_path=file_path,
                language=language,
                chunk_index=len(chunks),
                total_chunks=-1,
                start_line=current_start,
                end_line=current_start + current_text.count("\n"),
                content=current_text.strip(),
                content_hash=content_hash,
            ))
            current_start += current_text.count("\n") + 1
            current_text = section
        else:
            current_text = f"{current_text}\n\n{section}" if current_text else section

    if current_text.strip():
        chunks.append(Chunk(
            file_path=file_path,
            language=language,
            chunk_index=len(chunks),
            total_chunks=-1,
            start_line=current_start,
            end_line=current_start + current_text.count("\n"),
            content=current_text.strip(),
            content_hash=content_hash,
        ))

    for c in chunks:
        c.total_chunks = len(chunks)

    return chunks
